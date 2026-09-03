"""
Generates the official PakTradeX Hackathon Pitch Deck (.pptx)
"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # 16:9 widescreen format
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette
    BG_DARK = RGBColor(15, 23, 42)       # Slate 900
    BG_CARD = RGBColor(30, 41, 59)       # Slate 800
    PRIMARY = RGBColor(16, 185, 129)     # Emerald Green (PSX Bullish)
    PRIMARY_LIGHT = RGBColor(52, 211, 153)
    TEXT_MAIN = RGBColor(248, 250, 252)  # White/Slate 50
    TEXT_MUTED = RGBColor(148, 163, 184) # Slate 400
    ACCENT_BLUE = RGBColor(56, 189, 248) # Sky 400
    ACCENT_GOLD = RGBColor(251, 191, 36) # Amber 400

    def add_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background()
        return bg

    def add_header(slide, tag, title, subtitle):
        # Tag Badge
        tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.4))
        tf_tag = tag_box.text_frame
        tf_tag.word_wrap = True
        p_tag = tf_tag.paragraphs[0]
        p_tag.text = tag.upper()
        p_tag.font.size = Pt(11)
        p_tag.font.bold = True
        p_tag.font.color.rgb = PRIMARY

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_MAIN

        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.5))
        tf_sub = sub_box.text_frame
        tf_sub.word_wrap = True
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = subtitle
        p_sub.font.size = Pt(13)
        p_sub.font.color.rgb = TEXT_MUTED

    # ─────────────────────────────────────────────────────────────
    # SLIDE 1: Title & Cover
    # ─────────────────────────────────────────────────────────────
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide1)

    # Main Card
    card1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.2), Inches(10.9), Inches(5.1))
    card1.fill.solid()
    card1.fill.fore_color.rgb = BG_CARD
    card1.line.color.rgb = RGBColor(51, 65, 85)

    title_box = slide1.shapes.add_textbox(Inches(1.8), Inches(1.8), Inches(9.7), Inches(2.0))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    p0 = tf1.paragraphs[0]
    p0.text = "🇵🇰 PakTradeX"
    p0.font.size = Pt(44)
    p0.font.bold = True
    p0.font.color.rgb = PRIMARY_LIGHT

    p1 = tf1.add_paragraph()
    p1.text = "Pakistan's Premier Next-Gen PSX Trading & Investment Platform"
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_MAIN
    p1.space_before = Pt(12)

    p2 = tf1.add_paragraph()
    p2.text = "Democratizing retail stock market participation through real-time PSX data, risk-free simulation, and AI-powered intelligence."
    p2.font.size = Pt(13)
    p2.font.color.rgb = TEXT_MUTED
    p2.space_before = Pt(8)

    # Meta Badges
    badges_box = slide1.shapes.add_textbox(Inches(1.8), Inches(4.5), Inches(9.7), Inches(1.2))
    tf_b = badges_box.text_frame
    p_b = tf_b.paragraphs[0]
    p_b.text = "⚡ Full-Stack Flutter & FastAPI  •  📊 Real PSX Market Feed  •  🛡️ 1-Time KYC & Raast Rails  •  🤖 Gemini AI Copilot"
    p_b.font.size = Pt(12)
    p_b.font.bold = True
    p_b.font.color.rgb = ACCENT_BLUE

    # ─────────────────────────────────────────────────────────────
    # SLIDE 2: The Problem & Who It Affects
    # ─────────────────────────────────────────────────────────────
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide2)
    add_header(slide2, "1. Problem Discovery", "The Problem & Who It Affects", "Why less than 0.1% of Pakistan's 240M population invests in equities.")

    problems = [
        ("Massive Financial Exclusion", "Less than 250,000 active individual investor accounts exist in PSX out of 240M+ citizens due to high perceived risk and zero practical learning avenues.", ACCENT_GOLD),
        ("Intimidating & Outdated Tools", "Existing broker portals are clunky, desktop-bound, lack paper-trading sandboxes, and offer no modern mobile-first user experience.", ACCENT_BLUE),
        ("Complex Onboarding & Friction", "Lengthy physical documentation, absence of instant KYC verification, and lack of integration with modern digital rails (Raast, JazzCash).", PRIMARY),
    ]

    for i, (title, desc, color) in enumerate(problems):
        x = Inches(0.8 + (i * 3.95))
        box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.2), Inches(3.75), Inches(4.5))
        box.fill.solid()
        box.fill.fore_color.rgb = BG_CARD
        box.line.color.rgb = color

        tb = slide2.shapes.add_textbox(x + Inches(0.2), Inches(2.5), Inches(3.35), Inches(3.9))
        tf = tb.text_frame
        tf.word_wrap = True
        
        pt = tf.paragraphs[0]
        pt.text = f"0{i+1}. {title}"
        pt.font.size = Pt(16)
        pt.font.bold = True
        pt.font.color.rgb = color

        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(12)
        pd.font.color.rgb = TEXT_MAIN
        pd.space_before = Pt(14)

    # ─────────────────────────────────────────────────────────────
    # SLIDE 3: Solution & Audience
    # ─────────────────────────────────────────────────────────────
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide3)
    add_header(slide3, "2. Value Proposition", "Our Solution & Target Audience", "An institutional-grade trading terminal built for the next generation of Pakistani investors.")

    cards3 = [
        ("Interactive Paper Trading", "Risk-free virtual portfolio (1,000,000 PKR virtual balance) running on live PSX market prices with Market & Limit order execution.", "Aspiring Retail Investors & Students"),
        ("Shariah-Compliant Screening", "Dedicated KMI-30 tracking and automatic Islamic Shariah compliance tags for halal equity investing.", "Ethical & Shariah-Conscious Investors"),
        ("Instant Digital Rails & KYC", "Simulated SECP-compliant 1-time KYC (CNIC + Bank account binding) with instant Raast, IBFT, JazzCash & EasyPaisa payments.", "Digital-First Youth & Mobile Users"),
        ("Gemini AI Market Copilot", "On-demand company fundamental breakdown, live sentiment analysis, and risk rulebook summaries in plain Urdu & English.", "Everyday Traders seeking Alpha"),
    ]

    for i, (feat, aud, sub) in enumerate(cards3):
        r = i // 2
        c = i % 2
        x = Inches(0.8 + (c * 5.95))
        y = Inches(2.2 + (r * 2.45))
        
        box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.75), Inches(2.25))
        box.fill.solid()
        box.fill.fore_color.rgb = BG_CARD
        box.line.color.rgb = RGBColor(51, 65, 85)

        tb = slide3.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), Inches(5.35), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = f"✨ {feat}"
        p1.font.size = Pt(15)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY_LIGHT

        p2 = tf.add_paragraph()
        p2.text = aud
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MAIN
        p2.space_before = Pt(4)

        p3 = tf.add_paragraph()
        p3.text = f"Target: {sub}"
        p3.font.size = Pt(10)
        p3.font.bold = True
        p3.font.color.rgb = ACCENT_GOLD
        p3.space_before = Pt(6)

    # ─────────────────────────────────────────────────────────────
    # SLIDE 4: Need & Impact
    # ─────────────────────────────────────────────────────────────
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide4)
    add_header(slide4, "3. Market Impact", "The Need It Addresses & The Impact It Makes", "Empowering individuals, accelerating capital formation, and fostering financial literacy.")

    impacts = [
        ("Financial Literacy at Scale", "Converts passive savers battling inflation into informed equity investors by providing a zero-risk educational sandbox with real market dynamics.", "📈 Educational Empowerment"),
        ("Capital Market Inflow", "Prepares everyday citizens with hands-on trading experience so they transition confidently into real CDC-regulated broker accounts.", "💰 Economic Growth"),
        ("Inclusive Digital Finance", "Bridges the gap between modern fintech (Raast / Microfinance wallets) and traditional PSX equity brokers through frictionless UI.", "🌐 Financial Inclusion"),
    ]

    for i, (title, desc, badge) in enumerate(impacts):
        x = Inches(0.8 + (i * 3.95))
        box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.2), Inches(3.75), Inches(4.5))
        box.fill.solid()
        box.fill.fore_color.rgb = BG_CARD
        box.line.color.rgb = ACCENT_BLUE

        tb = slide4.shapes.add_textbox(x + Inches(0.2), Inches(2.5), Inches(3.35), Inches(3.9))
        tf = tb.text_frame
        tf.word_wrap = True

        pb = tf.paragraphs[0]
        pb.text = badge
        pb.font.size = Pt(11)
        pb.font.bold = True
        pb.font.color.rgb = ACCENT_GOLD

        pt = tf.add_paragraph()
        pt.text = title
        pt.font.size = Pt(16)
        pt.font.bold = True
        pt.font.color.rgb = PRIMARY_LIGHT
        pt.space_before = Pt(8)

        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(12)
        pd.font.color.rgb = TEXT_MAIN
        pd.space_before = Pt(12)

    # ─────────────────────────────────────────────────────────────
    # SLIDE 5: Innovation & Technology
    # ─────────────────────────────────────────────────────────────
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide5)
    add_header(slide5, "4. Technical Architecture", "Innovation & Technology Stack", "High-frequency reactive client architecture backed by asynchronous Python pipelines.")

    tech_cards = [
        ("Flutter 3.x + Riverpod", "Cross-platform client with reactive state management, sub-second price animations, and zero-downtime offline fallback engine.", ACCENT_BLUE),
        ("FastAPI + Async Python", "High-concurrency backend (Python 3.11) with non-blocking endpoints, WebSockets, and SQLite / PostgreSQL persistence.", PRIMARY),
        ("Live PSX Engine via Yahoo Finance", "Concurrent multi-threaded pipeline querying PSX (.KA tickers) with real 5-day OHLCV candles & 5-level order book depth.", ACCENT_GOLD),
        ("Google Gemini AI Copilot", "Integrated Generative AI financial assistant answering fundamental questions and generating daily KSE market briefs.", PRIMARY_LIGHT),
    ]

    for i, (title, desc, color) in enumerate(tech_cards):
        r = i // 2
        c = i % 2
        x = Inches(0.8 + (c * 5.95))
        y = Inches(2.2 + (r * 2.45))
        
        box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.75), Inches(2.25))
        box.fill.solid()
        box.fill.fore_color.rgb = BG_CARD
        box.line.color.rgb = color

        tb = slide5.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), Inches(5.35), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(15)
        p1.font.bold = True
        p1.font.color.rgb = color

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = TEXT_MAIN
        p2.space_before = Pt(6)

    # ─────────────────────────────────────────────────────────────
    # SLIDE 6: Feasibility & What We Actually Built
    # ─────────────────────────────────────────────────────────────
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide6)
    add_header(slide6, "5. Execution & Feasibility", "Feasibility & What We Have Actually Built", "Not a concept — a complete, working, production-tested mobile application.")

    built_items = [
        ("📱 Production Mobile App (APK)", "Built & compiled release APK (53.6 MB) with 100% working UI, navigation, and persistent storage.", PRIMARY),
        ("🔴 Live PSX Ticker Feed (31 Stocks)", "Real live prices for MCB (402.68), UBL (458.70), MEBL (574.34), SYS (130.82), OGDC (325.20), etc.", ACCENT_GOLD),
        ("⚡ Order Engine & Depth Charts", "Live Market/Limit buy & sell execution, order book depth, portfolio calculation, and transaction history.", ACCENT_BLUE),
        ("🛡️ SECP KYC & Raast Deposit Rails", "1-Time CNIC verification modal, Bank account binding, and Raast/JazzCash deposit & withdrawal simulation.", PRIMARY_LIGHT),
        ("🧪 100% Automated Test Coverage", "18/18 Unit & Widget test suites passed with automated GitHub Actions CI/CD release workflow.", PRIMARY),
    ]

    for i, (title, desc, color) in enumerate(built_items):
        y = Inches(2.1 + (i * 0.98))
        box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.73), Inches(0.85))
        box.fill.solid()
        box.fill.fore_color.rgb = BG_CARD
        box.line.color.rgb = color

        tb = slide6.shapes.add_textbox(Inches(1.0), y + Inches(0.08), Inches(11.3), Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"{title} — "
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = color

        run = p.add_run()
        run.text = desc
        run.font.size = Pt(11)
        run.font.bold = False
        run.font.color.rgb = TEXT_MAIN

    output_path = "PakTradeX_Hackathon_Pitch.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_presentation()
